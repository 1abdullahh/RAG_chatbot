"""
requirements.txt (add these):
streamlit>=1.40         # needed for chat_input(accept_file=...) "+" attach icon
langchain-core>=0.3
langchain-openai>=0.2
langchain-community>=0.3
langchain-text-splitters>=0.3
faiss-cpu>=1.8
pypdf>=4.2
openai>=1.40
python-docx>=1.1
python-pptx>=0.6.23
pandas>=2.2
openpyxl>=3.1
pytesseract>=0.3.10
pillow>=10.3

Note: pytesseract needs the Tesseract OCR engine installed on the system
(e.g. `apt-get install tesseract-ocr` on Linux, or `brew install tesseract` on Mac).
Without it, image files will fail to process.
"""

import io
import os
import tempfile
import uuid

import pandas as pd
import pytesseract
import streamlit as st
from docx import Document as DocxReader
from openai import OpenAI
from PIL import Image
from pptx import Presentation

from langchain_community.document_loaders import PyPDFLoader
from langchain_community.vectorstores import FAISS
from langchain_core.documents import Document
from langchain_core.output_parsers import StrOutputParser
from langchain_core.prompts import ChatPromptTemplate
from langchain_core.runnables import RunnablePassthrough
from langchain_openai import ChatOpenAI, OpenAIEmbeddings
from langchain_text_splitters import RecursiveCharacterTextSplitter

# ----------------------------------------
# Page setup
# ----------------------------------------
st.set_page_config(page_title="Multi-Format RAG Agent", page_icon="📄", layout="wide")

SUPPORTED_EXTENSIONS = ["pdf", "docx", "pptx", "csv", "xlsx", "xls", "png", "jpg", "jpeg"]

# ----------------------------------------
# STEP 1: API key gate
# Nothing else renders until a valid key is provided.
# ----------------------------------------
if "api_key_valid" not in st.session_state:
    st.session_state.api_key_valid = False
if "api_key" not in st.session_state:
    st.session_state.api_key = ""


def validate_api_key(key: str) -> tuple[bool, str]:
    """Ping OpenAI to check whether the key works, without doing a billed completion."""
    if not key or not key.strip():
        return False, "Please enter an API key."
    try:
        client = OpenAI(api_key=key)
        client.models.list()  # cheap call just to confirm the key is accepted
        return True, ""
    except Exception as e:
        return False, str(e)


if not st.session_state.api_key_valid:
    st.title("🔑 Enter Your OpenAI API Key")
    st.caption("Your key is only kept for this session and is never stored on a server.")

    with st.form("api_key_form"):
        key_input = st.text_input("OpenAI API Key", type="password")
        submitted = st.form_submit_button("Validate & Continue")

    if submitted:
        with st.spinner("Validating key…"):
            is_valid, error_msg = validate_api_key(key_input)
        if is_valid:
            st.session_state.api_key_valid = True
            st.session_state.api_key = key_input
            os.environ["OPENAI_API_KEY"] = key_input
            st.rerun()
        else:
            st.error(f"❌ Invalid API key.")

    st.stop()  # don't render anything below until the key is valid

os.environ["OPENAI_API_KEY"] = st.session_state.api_key

# ----------------------------------------
# STEP 2: Chat session storage
# Every chat (its messages + its own indexed files) lives in st.session_state.chats,
# keyed by a chat_id. Switching / creating chats just changes current_chat_id,
# so old chats are never lost — that's the "history".
# ----------------------------------------
if "chats" not in st.session_state:
    st.session_state.chats = {}
if "current_chat_id" not in st.session_state:
    st.session_state.current_chat_id = None


def new_chat() -> str:
    chat_id = str(uuid.uuid4())
    st.session_state.chats[chat_id] = {
        "title": "New chat",
        "messages": [],
        "files": {},          # filename -> bytes
        "vectorstore": None,
        "indexed_sig": None,  # tuple of (filenames, chunk_size, chunk_overlap) already indexed
    }
    st.session_state.current_chat_id = chat_id
    return chat_id


if not st.session_state.chats:
    new_chat()
elif st.session_state.current_chat_id not in st.session_state.chats:
    st.session_state.current_chat_id = next(iter(st.session_state.chats))

current_chat = st.session_state.chats[st.session_state.current_chat_id]

# ----------------------------------------
# Sidebar: settings + new chat + chat history
# ----------------------------------------
with st.sidebar:
    st.header("Settings")
    if st.button("🔓 Change API Key"):
        st.session_state.api_key_valid = False
        st.session_state.api_key = ""
        st.rerun()

    model_name = st.selectbox("Chat model", ["gpt-4o-mini", "gpt-4o", "gpt-4.1-mini"])
    chunk_size = st.slider("Chunk size", 500, 2000, 1000, step=100)
    chunk_overlap = st.slider("Chunk overlap", 0, 400, 150, step=50)
    top_k = st.slider("Retrieved chunks (k)", 2, 10, 4)
    st.caption("ℹ️ Image OCR requires the Tesseract engine installed on the host machine.")

    st.divider()

    if st.button("➕ New Chat", use_container_width=True):
        new_chat()
        st.rerun()

    st.subheader("Chat history")
    # Most recently created chats first
    for chat_id in reversed(list(st.session_state.chats.keys())):
        chat = st.session_state.chats[chat_id]
        label = chat["title"]
        is_current = chat_id == st.session_state.current_chat_id
        if st.button(
            ("• " if is_current else "") + label,
            key=f"hist_{chat_id}",
            use_container_width=True,
            type="primary" if is_current else "secondary",
        ):
            st.session_state.current_chat_id = chat_id
            st.rerun()

# ----------------------------------------
# Main header
# ----------------------------------------
st.title("📄 Multi-Format RAG Agent")
st.caption("Attach PDF, Word, PowerPoint, CSV, Excel, or image files with the **+** button below — ask questions grounded in their content.")

# ----------------------------------------
# File loaders for each supported type
# ----------------------------------------

def load_pdf(file_bytes: bytes, filename: str) -> list[Document]:
    with tempfile.NamedTemporaryFile(delete=False, suffix=".pdf") as tmp:
        tmp.write(file_bytes)
        tmp_path = tmp.name
    try:
        docs = PyPDFLoader(tmp_path).load()
    finally:
        os.unlink(tmp_path)
    for d in docs:
        d.metadata["source"] = filename
        d.metadata["location"] = f"page {d.metadata.get('page', '?')}"
    return docs


def load_docx(file_bytes: bytes, filename: str) -> list[Document]:
    doc = DocxReader(io.BytesIO(file_bytes))
    text = "\n".join(p.text for p in doc.paragraphs if p.text.strip())
    return [Document(page_content=text, metadata={"source": filename, "location": "document"})]


def load_pptx(file_bytes: bytes, filename: str) -> list[Document]:
    prs = Presentation(io.BytesIO(file_bytes))
    docs = []
    for i, slide in enumerate(prs.slides, start=1):
        texts = []
        for shape in slide.shapes:
            if shape.has_text_frame:
                texts.append(shape.text_frame.text)
        content = "\n".join(t for t in texts if t.strip())
        if content.strip():
            docs.append(Document(page_content=content, metadata={"source": filename, "location": f"slide {i}"}))
    return docs


def load_csv(file_bytes: bytes, filename: str) -> list[Document]:
    df = pd.read_csv(io.BytesIO(file_bytes))
    docs = []
    for i, row in df.iterrows():
        content = ", ".join(f"{col}: {row[col]}" for col in df.columns)
        docs.append(Document(page_content=content, metadata={"source": filename, "location": f"row {i + 1}"}))
    return docs


def load_excel(file_bytes: bytes, filename: str) -> list[Document]:
    sheets = pd.read_excel(io.BytesIO(file_bytes), sheet_name=None)
    docs = []
    for sheet_name, df in sheets.items():
        content = df.to_string(index=False)
        docs.append(Document(page_content=content, metadata={"source": filename, "location": f"sheet '{sheet_name}'"}))
    return docs


def load_image(file_bytes: bytes, filename: str) -> list[Document]:
    image = Image.open(io.BytesIO(file_bytes))
    text = pytesseract.image_to_string(image)
    if not text.strip():
        text = "[No readable text found in this image.]"
    return [Document(page_content=text, metadata={"source": filename, "location": "image (OCR)"})]


LOADERS = {
    "pdf": load_pdf,
    "docx": load_docx,
    "pptx": load_pptx,
    "csv": load_csv,
    "xlsx": load_excel,
    "xls": load_excel,
    "png": load_image,
    "jpg": load_image,
    "jpeg": load_image,
}


def load_bytes(file_bytes: bytes, filename: str) -> list[Document]:
    ext = filename.rsplit(".", 1)[-1].lower()
    loader_fn = LOADERS.get(ext)
    if loader_fn is None:
        st.warning(f"Skipped unsupported file: {filename}")
        return []
    try:
        return loader_fn(file_bytes, filename)
    except Exception as e:
        st.warning(f"Could not read {filename}: {e}")
        return []


def build_vectorstore(files: dict, size: int, overlap: int) -> FAISS:
    all_docs = []
    for filename, file_bytes in files.items():
        all_docs.extend(load_bytes(file_bytes, filename))

    if not all_docs:
        raise ValueError("No readable content found in the uploaded files.")

    splitter = RecursiveCharacterTextSplitter(
        chunk_size=size,
        chunk_overlap=overlap,
        separators=["\n\n", "\n", ". ", " ", ""],
    )
    chunks = splitter.split_documents(all_docs)

    embeddings = OpenAIEmbeddings(model="text-embedding-3-small")
    return FAISS.from_documents(chunks, embeddings)


def ensure_indexed(chat: dict, size: int, overlap: int):
    """(Re)build the chat's vectorstore if its files or chunk settings changed."""
    sig = (tuple(sorted(chat["files"].keys())), size, overlap)
    if not chat["files"]:
        return
    if chat["indexed_sig"] != sig:
        with st.spinner("Reading and indexing files…"):
            try:
                chat["vectorstore"] = build_vectorstore(chat["files"], size, overlap)
                chat["indexed_sig"] = sig
            except ValueError as e:
                st.error(str(e))

# ----------------------------------------
# RAG chain
# ----------------------------------------
RAG_PROMPT = ChatPromptTemplate.from_messages([
    ("system",
     "You are a helpful assistant that answers questions strictly using the "
     "provided context extracted from the user's uploaded files.\n"
     "Rules:\n"
     "1. Answer ONLY from the context below.\n"
     "2. If the answer is not in the context, say: "
     "\"I couldn't find that in the document.\"\n"
     "3. Cite the source file and location you used, e.g. (report.pdf, page 3).\n\n"
     "Context:\n{context}"),
    ("human", "{question}"),
])


def format_docs(docs) -> str:
    return "\n\n".join(
        f"[{d.metadata.get('source', 'unknown')} — {d.metadata.get('location', '?')}] {d.page_content}"
        for d in docs
    )


def get_chain(vectorstore: FAISS, model: str, k: int):
    retriever = vectorstore.as_retriever(search_kwargs={"k": k})
    llm = ChatOpenAI(model=model, temperature=0)
    return (
        {"context": retriever | format_docs, "question": RunnablePassthrough()}
        | RAG_PROMPT
        | llm
        | StrOutputParser()
    ), retriever

# ----------------------------------------
# Chat interface
# ----------------------------------------
for msg in current_chat["messages"]:
    with st.chat_message(msg["role"]):
        if msg.get("attachments"):
            st.caption("📎 " + ", ".join(msg["attachments"]))
        st.markdown(msg["content"])

if current_chat["files"]:
    st.caption(f"📁 Indexed in this chat: {', '.join(current_chat['files'].keys())}")

# The "+" icon that opens the OS file browser lives inside chat_input itself.
chat_value = st.chat_input(
    "Ask a question, or attach files with +…",
    accept_file="multiple",
    file_type=SUPPORTED_EXTENSIONS,
)

if chat_value:
    question = chat_value.text
    new_files = chat_value.files or []

    # Rename the chat the first time something is sent
    if current_chat["title"] == "New chat" and question:
        current_chat["title"] = question[:40] + ("…" if len(question) > 40 else "")

    attachment_names = []
    for f in new_files:
        current_chat["files"][f.name] = f.getvalue()
        attachment_names.append(f.name)

    if attachment_names:
        ensure_indexed(current_chat, chunk_size, chunk_overlap)

    if question:
        current_chat["messages"].append(
            {"role": "user", "content": question, "attachments": attachment_names}
        )
        with st.chat_message("user"):
            if attachment_names:
                st.caption("📎 " + ", ".join(attachment_names))
            st.markdown(question)

        with st.chat_message("assistant"):
            if current_chat["vectorstore"] is None:
                answer = "Please attach at least one file (using the **+** button) so I have something to answer from."
                st.markdown(answer)
            else:
                ensure_indexed(current_chat, chunk_size, chunk_overlap)
                chain, retriever = get_chain(current_chat["vectorstore"], model_name, top_k)
                with st.spinner("Thinking…"):
                    answer = chain.invoke(question)
                    st.markdown(answer)
                    with st.expander("🔍 Sources (retrieved chunks)"):
                        for doc in retriever.invoke(question):
                            st.markdown(f"**{doc.metadata.get('source', 'unknown')} — {doc.metadata.get('location', '?')}**")
                            st.text(doc.page_content[:500])
                            st.divider()

        current_chat["messages"].append({"role": "assistant", "content": answer, "attachments": []})
    elif attachment_names:
        st.toast(f"Indexed: {', '.join(attachment_names)}")

    st.rerun()
elif not current_chat["messages"] and not current_chat["files"]:
    st.info("")